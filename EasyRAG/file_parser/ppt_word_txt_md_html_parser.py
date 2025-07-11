import os
import tempfile
import json
import mimetypes
from typing import Any, Dict
from EasyRAG.file_parser.document_parser import DocumentParser


class PPTWordTxtMDHTMLParser(DocumentParser):
    """
    解析ppt、word、txt、md、html文件，支持图片和影印文本提取
    """
    def __init__(self):
        super().__init__()
        self._setup_magic_pdf_config()
    
    def _setup_magic_pdf_config(self):
        """设置 magic_pdf 配置文件"""
        try:
            # 创建 magic_pdf 配置文件
            config_dir = os.path.expanduser("~")
            config_file = os.path.join(config_dir, "magic-pdf.json")
            
            if not os.path.exists(config_file):
                # 创建默认配置，优化多格式文档支持
                default_config = {
                    "latex_delimiters": {
                        "inline": ["$", "$"],
                        "display": ["$$", "$$"]
                    },
                    "output_format": "markdown",
                    "enable_ocr": True,
                    "enable_formula": True,
                    "latex-delimiter-config": {
                        "inline": ["$", "$"],
                        "display": ["$$", "$$"]
                    },
                    # OCR 相关配置
                    "ocr_config": {
                        "lang": "ch",  # 中文优先
                        "detect_orientation": True,
                        "preprocessing": True,
                        "confidence_threshold": 0.5
                    },
                    # 图片处理配置
                    "image_config": {
                        "extract_images": True,
                        "image_quality": "high",
                        "image_format": "png"
                    },
                    # 多格式文档支持
                    "document_config": {
                        "extract_tables": True,
                        "extract_images": True,
                        "extract_equations": True,
                        "preserve_formatting": True
                    }
                }
                
                with open(config_file, 'w', encoding='utf-8') as f:
                    json.dump(default_config, f, indent=2, ensure_ascii=False)
                    
        except Exception as e:
            print(f"警告: 无法创建 magic_pdf 配置文件: {e}")
    
    def parse(self, doc_info: Dict[str, Any], 
              file_info: Dict[str, Any], 
              knowledge_base_info: Dict[str, Any], 
              config: Dict[str, Any]):
        """
        解析多种格式的文档文件
        
        Args:
            doc_info: 文档信息
            file_info: 文件信息，包含 file_content
            knowledge_base_info: 知识库信息
            config: 解析配置
            
        Returns:
            解析结果
        """
        try:
            # 获取文件扩展名
            file_name = doc_info.get('file_name', 'unknown')
            file_extension = self._get_file_extension(file_name)
            
            # 创建临时文件
            temp_dir = tempfile.gettempdir()
            temp_file_path = os.path.join(temp_dir, f"{doc_info.get('doc_id', 'temp')}{file_extension}")
            
            # 写入文件内容
            with open(temp_file_path, "wb") as f:
                f.write(file_info['file_content'])
            
            # 根据文件类型选择合适的解析方法
            result = self._parse_by_file_type(temp_file_path, file_extension, config)
            
            # 清理临时文件
            try:
                os.remove(temp_file_path)
            except:
                pass
            
            return result
            
        except Exception as e:
            print(f"文档解析失败: {e}")
            return {
                'success': False,
                'error': str(e),
                'content': '',
                'metadata': {}
            }
    
    def _get_file_extension(self, file_name: str) -> str:
        """获取文件扩展名"""
        if '.' in file_name:
            return '.' + file_name.split('.')[-1].lower()
        return ''
    
    def _parse_by_file_type(self, file_path: str, file_extension: str, config: Dict[str, Any]) -> Dict[str, Any]:
        """
        根据文件类型选择合适的解析方法
        
        Args:
            file_path: 文件路径
            file_extension: 文件扩展名
            config: 解析配置
            
        Returns:
            解析结果
        """
        file_extension = file_extension.lower()
        
        if file_extension in ['.ppt', '.pptx']:
            return self._parse_presentation(file_path, config)
        elif file_extension in ['.doc', '.docx']:
            return self._parse_word_document(file_path, config)
        elif file_extension in ['.txt']:
            return self._parse_text_file(file_path, config)
        elif file_extension in ['.md', '.markdown']:
            return self._parse_markdown_file(file_path, config)
        elif file_extension in ['.html', '.htm']:
            return self._parse_html_file(file_path, config)
        else:
            # 尝试使用通用方法
            return self._parse_with_magic_pdf_generic(file_path, config)
    
    def _parse_presentation(self, file_path: str, config: Dict[str, Any]) -> Dict[str, Any]:
        """
        解析PPT文件
        
        Args:
            file_path: PPT文件路径
            config: 解析配置
            
        Returns:
            解析结果
        """
        try:
            # 使用 python-pptx 解析PPT
            from pptx import Presentation
            
            prs = Presentation(file_path)
            content = ""
            slide_count = len(prs.slides)
            
            for i, slide in enumerate(prs.slides):
                content += f"\n--- 第 {i + 1} 页 ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content += shape.text.strip() + "\n"
                    
                    # 处理图片
                    if shape.shape_type == 13:  # 图片类型
                        content += "[图片]\n"
                
                content += "\n"
            
            return {
                'success': True,
                'content': content,
                'metadata': {
                    'file_type': 'presentation',
                    'slides': slide_count,
                    'format': 'text',
                    'parser': 'python-pptx'
                }
            }
            
        except Exception as e:
            print(f"PPT解析失败，尝试使用magic_pdf: {e}")
            return self._parse_with_magic_pdf_generic(file_path, config)
    
    def _parse_word_document(self, file_path: str, config: Dict[str, Any]) -> Dict[str, Any]:
        """
        解析Word文档
        
        Args:
            file_path: Word文件路径
            config: 解析配置
            
        Returns:
            解析结果
        """
        try:
            # 使用 python-docx 解析Word文档
            from docx import Document
            
            doc = Document(file_path)
            content = ""
            
            # 提取段落文本
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content += paragraph.text.strip() + "\n"
            
            # 提取表格
            for table in doc.tables:
                content += "\n[表格]\n"
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    content += row_text + "\n"
                content += "\n"
            
            return {
                'success': True,
                'content': content,
                'metadata': {
                    'file_type': 'word_document',
                    'paragraphs': len(doc.paragraphs),
                    'tables': len(doc.tables),
                    'format': 'text',
                    'parser': 'python-docx'
                }
            }
            
        except Exception as e:
            print(f"Word文档解析失败，尝试使用magic_pdf: {e}")
            return self._parse_with_magic_pdf_generic(file_path, config)
    
    def _parse_text_file(self, file_path: str, config: Dict[str, Any]) -> Dict[str, Any]:
        """
        解析文本文件
        
        Args:
            file_path: 文本文件路径
            config: 解析配置
            
        Returns:
            解析结果
        """
        try:
            # 尝试多种编码
            encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
            content = ""
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    break
                except UnicodeDecodeError:
                    continue
            
            if not content:
                # 如果文本解析失败，尝试使用magic_pdf处理可能的影印文本
                return self._parse_with_magic_pdf_generic(file_path, config)
            
            return {
                'success': True,
                'content': content,
                'metadata': {
                    'file_type': 'text',
                    'encoding': encoding,
                    'format': 'text',
                    'parser': 'native'
                }
            }
            
        except Exception as e:
            print(f"文本文件解析失败: {e}")
            return {
                'success': False,
                'error': str(e),
                'content': '',
                'metadata': {}
            }
    
    def _parse_markdown_file(self, file_path: str, config: Dict[str, Any]) -> Dict[str, Any]:
        """
        解析Markdown文件
        
        Args:
            file_path: Markdown文件路径
            config: 解析配置
            
        Returns:
            解析结果
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            return {
                'success': True,
                'content': content,
                'metadata': {
                    'file_type': 'markdown',
                    'format': 'markdown',
                    'parser': 'native'
                }
            }
            
        except Exception as e:
            print(f"Markdown文件解析失败: {e}")
            return {
                'success': False,
                'error': str(e),
                'content': '',
                'metadata': {}
            }
    
    def _parse_html_file(self, file_path: str, config: Dict[str, Any]) -> Dict[str, Any]:
        """
        解析HTML文件
        
        Args:
            file_path: HTML文件路径
            config: 解析配置
            
        Returns:
            解析结果
        """
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            # 使用BeautifulSoup解析HTML
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # 提取文本内容
            text_content = soup.get_text(separator='\n', strip=True)
            
            # 提取图片信息
            images = soup.find_all('img')
            image_info = []
            for img in images:
                src = img.get('src', '')
                alt = img.get('alt', '')
                if src:
                    image_info.append(f"[图片: {alt or src}]")
            
            # 组合内容
            content = text_content
            if image_info:
                content += "\n\n" + "\n".join(image_info)
            
            return {
                'success': True,
                'content': content,
                'metadata': {
                    'file_type': 'html',
                    'images': len(images),
                    'format': 'text',
                    'parser': 'beautifulsoup'
                }
            }
            
        except Exception as e:
            print(f"HTML文件解析失败: {e}")
            return {
                'success': False,
                'error': str(e),
                'content': '',
                'metadata': {}
            }
    
    def _parse_with_magic_pdf_generic(self, file_path: str, config: Dict[str, Any]) -> Dict[str, Any]:
        """
        使用magic_pdf通用方法解析文件，支持图片和影印文本
        
        Args:
            file_path: 文件路径
            config: 解析配置
            
        Returns:
            解析结果
        """
        try:
            # 导入magic_pdf相关模块
            from magic_pdf.pdf_parse_union_core_v2 import pdf_parse_union
            from magic_pdf.data.data_reader_writer import FileBasedDataReader, FileBasedDataWriter
            from magic_pdf.data.dataset import PymuDocDataset
            from magic_pdf.config.enums import SupportedPdfParseMethod
            
            # 读取文件
            reader = FileBasedDataReader("")
            file_bytes = reader.read(file_path)
            ds = PymuDocDataset(file_bytes)
            
            # 自动检测文件类型
            parse_method = ds.classify()
            is_ocr_needed = parse_method == SupportedPdfParseMethod.OCR
            
            print(f"文件类型检测结果: {parse_method.value} ({'需要OCR' if is_ocr_needed else '文本模式'})")
            
            # 设置输出目录
            temp_dir = tempfile.gettempdir()
            output_dir = os.path.join(temp_dir, f"magic_pdf_output_{os.path.basename(file_path)}")
            os.makedirs(output_dir, exist_ok=True)
            
            # 创建图片写入器
            image_writer = FileBasedDataWriter(output_dir)
            
            # 调用magic_pdf解析
            result = pdf_parse_union(
                model_list=[],  # 使用默认模型
                dataset=ds,
                imageWriter=image_writer,
                parse_mode=parse_method.value,
                debug_mode=config.get('debug_mode', False),
                lang=config.get('ocr_lang', 'ch')  # 默认中文
            )
            
            # 提取解析结果
            try:
                content_list = result.get_content_list(os.path.basename(file_path))
            except Exception as e:
                print(f"获取内容列表失败: {e}")
                content_list = []
            
            # 处理内容列表
            extracted_text = self._extract_text_from_content_list(content_list)
            
            return {
                'success': True,
                'content': extracted_text,
                'metadata': {
                    'parse_method': parse_method.value,
                    'ocr_enabled': is_ocr_needed,
                    'pages': len(content_list) if content_list else 0,
                    'format': config.get('output_format', 'markdown'),
                    'images_extracted': len([c for c in content_list if c.get('type') == 'image']) if content_list else 0,
                    'tables_extracted': len([c for c in content_list if c.get('type') == 'table']) if content_list else 0,
                    'equations_extracted': len([c for c in content_list if c.get('type') == 'equation']) if content_list else 0,
                    'parser': 'magic_pdf'
                }
            }
            
        except Exception as e:
            print(f"magic_pdf通用解析失败: {e}")
            return {
                'success': False,
                'error': str(e),
                'content': '',
                'metadata': {}
            }
    
    def _extract_text_from_content_list(self, content_list: list) -> str:
        """
        从内容列表中提取文本
        
        Args:
            content_list: 内容列表
            
        Returns:
            提取的文本
        """
        if not content_list:
            return ""
        
        extracted_text = ""
        
        for item in content_list:
            item_type = item.get('type', '')
            
            if item_type == 'text':
                text = item.get('text', '').strip()
                if text:
                    extracted_text += text + "\n\n"
                    
            elif item_type == 'table':
                # 处理表格
                caption = item.get('table_caption', '')
                body = item.get('table_body', '')
                
                if caption:
                    extracted_text += f"表格标题: {caption}\n"
                if body:
                    extracted_text += f"表格内容:\n{body}\n\n"
                    
            elif item_type == 'equation':
                # 处理公式
                equation = item.get('text', '').strip()
                if equation:
                    extracted_text += f"公式: {equation}\n\n"
                    
            elif item_type == 'image':
                # 处理图片描述
                img_path = item.get('img_path', '')
                if img_path:
                    extracted_text += f"[图片: {img_path}]\n\n"
        
        return extracted_text.strip()
    
    def get_step_status(self, doc_id: str, step: str) -> Dict[str, Any]:
        """
        获取指定步骤的状态
        
        Args:
            doc_id: 文档ID
            step: 步骤名称
            
        Returns:
            Dict[str, Any]: 步骤状态信息
        """
        return {
            'doc_id': doc_id,
            'step': step,
            'status': 'completed',
            'message': '文档解析步骤已完成'
        }